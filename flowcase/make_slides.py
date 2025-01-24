import datetime
from typing import Dict, List

from pptx import Presentation  # as PP
from pptx.dml.color import RGBColor

from flowcase.reports import get_skills_keyword
from flowcase.types.cv import (
    Certification,
    Course,
    HonorsAward,
    ProjectExperienceExpanded,
)
from flowcase.types.cv import (
    Presentation as CVPresentation,
)


def add_experience_slide(
    prs: Presentation, layout, name: str, projects: List[ProjectExperienceExpanded]
):
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = f"{name} 👷🏻‍♀️"

    # text_frame to place text in
    tf = slide.shapes.placeholders[1].text_frame
    for pro in projects:
        p = tf.add_paragraph()
        p.text = f"{pro.customer.no}: {pro.description.no}"
        p.level = 0
        p = tf.add_paragraph()
        p.text = f'{", ".join(get_skills_keyword(pro.project_experience_skills))}'
        p.level = 1
    return slide


def add_course_slide(
    prs: Presentation,
    layout,
    name: str,
    courses: List[Course],
    background_color=(250, 233, 150),
):
    slide = prs.slides.add_slide(layout)
    # set background color
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*background_color)

    slide.shapes.title.text = f"Kurs: {name} 👩‍🎓"

    tf = slide.shapes.placeholders[1].text_frame
    for course in courses:
        p = tf.add_paragraph()
        p.text = f"{course.name.no}"
        p.level = 0
        p = tf.add_paragraph()
        p.text = f"{course.long_description.no}"
        p.level = 1
    return slide


def add_certification_slide(
    prs: Presentation,
    layout,
    name: str,
    certifications: List[Certification],
    background_color=(160, 250, 150),
):
    slide = prs.slides.add_slide(layout)

    # set background color
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*background_color)

    slide.shapes.title.text = f"Sertifisering: {name} 📜"

    # text_frame to place text in
    tf = slide.shapes.placeholders[1].text_frame
    for cert in certifications:
        p = tf.add_paragraph()
        p.text = f"{cert.name.no}"
        p.level = 0
        if cert.long_description.no:
            p = tf.add_paragraph()
            p.text = f"{cert.long_description.no}"
            p.level = 1
    return slide


def add_presentation_slide(
    prs: Presentation,
    layout,
    name: str,
    presentations: List[CVPresentation],
    background_color=(150, 250, 237),
):
    slide = prs.slides.add_slide(layout)

    # set background color
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*background_color)

    slide.shapes.title.text = f"Presentasjon: {name} 🙋🏻‍♂️"

    # text_frame to place text in
    tf = slide.shapes.placeholders[1].text_frame
    for pres in presentations:
        p = tf.add_paragraph()
        p.text = f"{pres.description.no}"
        p.level = 0
        if pres.long_description.no:
            p = tf.add_paragraph()
            p.text = f"{pres.description.no}"
            p.level = 1
    return slide


def add_honors_and_awards_slide(
    prs: Presentation,
    layout,
    name: str,
    honors_and_awards: List[HonorsAward],
    background_color=(250, 150, 237),
) -> bool:
    slide = prs.slides.add_slide(layout)

    # set background color
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*background_color)

    slide.shapes.title.text = f"Priser og utmerkslser: {name} 🏆"

    # text_frame to place text in
    tf = slide.shapes.placeholders[1].text_frame
    for hon in honors_and_awards:
        p = tf.add_paragraph()
        p.text = f"{hon.issuer.no} for {hon.for_work.no}"
        p.level = 0
        if hon.long_description.no:
            p = tf.add_paragraph()
            p.text = f"{hon.long_description.no}"
            p.level = 1
    return slide


def make_ppt_from_year_in_review(
    projects_worked_on: Dict[str, list[ProjectExperienceExpanded]],
    new_courses: Dict[str, list[Course]],
    new_certifications: Dict[str, list[Certification]],
    new_presentations: Dict[str, list[CVPresentation]],
    new_honors_and_awards: Dict[str, list[HonorsAward]],
    department_name: str = "Data engineering",
) -> bool:
    try:
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        bullet_slide_layout = prs.slide_layouts[1]

        title.text = f"{department_name} last year in review"
        subtitle.text = "Noa Ignite 🔥!"

        all_names = list(
            set(
                projects_worked_on.keys()
                | new_courses.keys()
                | new_certifications.keys()
                | new_presentations.keys()
                | new_honors_and_awards.keys()
            )
        )

        # loop over all people
        for name in all_names:
            if name in projects_worked_on:
                add_experience_slide(
                    prs, bullet_slide_layout, name, projects_worked_on[name]
                )

            if name in new_courses:
                add_course_slide(prs, bullet_slide_layout, name, new_courses[name])

            if name in new_certifications:
                add_certification_slide(
                    prs, bullet_slide_layout, name, new_certifications[name]
                )

            if name in new_presentations:
                add_presentation_slide(
                    prs, bullet_slide_layout, name, new_presentations[name]
                )

            if name in new_honors_and_awards:
                add_honors_and_awards_slide(
                    prs, bullet_slide_layout, name, new_honors_and_awards[name]
                )

        prs.save(
            f"presentation_{department_name}_{datetime.datetime.now().hour}{datetime.datetime.now().minute}.pptx"
        )
        return True
    except Exception as e:
        print(e)
        return False
